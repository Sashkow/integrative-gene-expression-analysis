#!/usr/bin/env python3

"""
Create PowerPoint Presentation with PCA Figures

Creates a PPTX with the combined PCA figure and editable text captions

@author Expression Integration Pipeline
@date 2025-11-15
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

print("\n=== Creating PowerPoint Presentation ===\n")

# Define paths
data_dir = "output/test_first_datasets/GSE55439_GSE93520_GSE28551_GSE100051"
pcas_dir = os.path.join(data_dir, "pcas")

# Image paths (we'll use the individual PCA plots)
# We need to generate individual panel images first
# For now, let's use the combined image and add text boxes

combined_image = os.path.join(pcas_dir, "combined_pca_figure.png")
output_file = os.path.join(pcas_dir, "combined_pca_figure.pptx")

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Add blank slide
blank_slide_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(blank_slide_layout)

# Add title
title_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)
)
title_frame = title_box.text_frame
title_frame.text = "PCA Analysis: Batch Effect Removal Assessment"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(20)
title_para.font.bold = True
title_para.alignment = PP_ALIGN.CENTER

# Add combined image
# Image should be about 8.27 x 5.85 inches (half A4)
img_width = Inches(8.27)
img_height = Inches(5.85)
img_left = Inches((10 - 8.27) / 2)  # Center horizontally
img_top = Inches(0.8)

slide.shapes.add_picture(
    combined_image,
    img_left, img_top,
    width=img_width, height=img_height
)

# Add panel labels as editable text boxes
# Position them over the image where A, B, C, D should be
label_size = Pt(18)
label_width = Inches(0.4)
label_height = Inches(0.4)

# Calculate positions for labels (top-left of each quadrant)
# The image is divided into 2x2 grid
panel_width = img_width / 2
panel_height = img_height / 2

# Panel A (top-left)
label_a = slide.shapes.add_textbox(
    img_left + Inches(0.1),
    img_top + Inches(0.1),
    label_width, label_height
)
label_a_frame = label_a.text_frame
label_a_frame.text = "A"
label_a_para = label_a_frame.paragraphs[0]
label_a_para.font.size = label_size
label_a_para.font.bold = True
label_a_para.font.color.rgb = RGBColor(0, 0, 0)

# Panel B (top-right)
label_b = slide.shapes.add_textbox(
    img_left + panel_width + Inches(0.1),
    img_top + Inches(0.1),
    label_width, label_height
)
label_b_frame = label_b.text_frame
label_b_frame.text = "B"
label_b_para = label_b_frame.paragraphs[0]
label_b_para.font.size = label_size
label_b_para.font.bold = True
label_b_para.font.color.rgb = RGBColor(0, 0, 0)

# Panel C (bottom-left)
label_c = slide.shapes.add_textbox(
    img_left + Inches(0.1),
    img_top + panel_height + Inches(0.1),
    label_width, label_height
)
label_c_frame = label_c.text_frame
label_c_frame.text = "C"
label_c_para = label_c_frame.paragraphs[0]
label_c_para.font.size = label_size
label_c_para.font.bold = True
label_c_para.font.color.rgb = RGBColor(0, 0, 0)

# Panel D (bottom-right)
label_d = slide.shapes.add_textbox(
    img_left + panel_width + Inches(0.1),
    img_top + panel_height + Inches(0.1),
    label_width, label_height
)
label_d_frame = label_d.text_frame
label_d_frame.text = "D"
label_d_para = label_d_frame.paragraphs[0]
label_d_para.font.size = label_size
label_d_para.font.bold = True
label_d_para.font.color.rgb = RGBColor(0, 0, 0)

# Add figure caption as editable text box at the bottom
caption_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(6.8), Inches(9), Inches(0.6)
)
caption_frame = caption_box.text_frame
caption_frame.word_wrap = True

caption_text = (
    "Figure 1. Principal component analysis comparing before and after batch effect removal. "
    "Panel A shows dataset distribution before ComBat correction, with strong technical variation. "
    "Panel B demonstrates successful batch correction with overlapping datasets. "
    "Panel C shows obscured trimester separation before correction. "
    "Panel D reveals clear biological separation by trimester after batch effect removal."
)

caption_frame.text = caption_text
caption_para = caption_frame.paragraphs[0]
caption_para.font.size = Pt(10)
caption_para.alignment = PP_ALIGN.LEFT

# Add panel descriptions as editable text boxes
desc_top = img_top + img_height + Inches(0.05)
desc_font_size = Pt(9)
desc_box_height = Inches(0.25)

# Panel A description
desc_a = slide.shapes.add_textbox(
    img_left, desc_top,
    panel_width - Inches(0.05), desc_box_height
)
desc_a_frame = desc_a.text_frame
desc_a_frame.text = "A: Before ComBat - by Dataset"
desc_a_para = desc_a_frame.paragraphs[0]
desc_a_para.font.size = desc_font_size
desc_a_para.alignment = PP_ALIGN.CENTER

# Panel B description
desc_b = slide.shapes.add_textbox(
    img_left + panel_width + Inches(0.05), desc_top,
    panel_width - Inches(0.05), desc_box_height
)
desc_b_frame = desc_b.text_frame
desc_b_frame.text = "B: After ComBat - by Dataset"
desc_b_para = desc_b_frame.paragraphs[0]
desc_b_para.font.size = desc_font_size
desc_b_para.alignment = PP_ALIGN.CENTER

# Panel C and D descriptions
desc_bottom = desc_top + Inches(0.3)

desc_c = slide.shapes.add_textbox(
    img_left, desc_bottom,
    panel_width - Inches(0.05), desc_box_height
)
desc_c_frame = desc_c.text_frame
desc_c_frame.text = "C: Before ComBat - by Trimester"
desc_c_para = desc_c_frame.paragraphs[0]
desc_c_para.font.size = desc_font_size
desc_c_para.alignment = PP_ALIGN.CENTER

desc_d = slide.shapes.add_textbox(
    img_left + panel_width + Inches(0.05), desc_bottom,
    panel_width - Inches(0.05), desc_box_height
)
desc_d_frame = desc_d.text_frame
desc_d_frame.text = "D: After ComBat - by Trimester"
desc_d_para = desc_d_frame.paragraphs[0]
desc_d_para.font.size = desc_font_size
desc_d_para.alignment = PP_ALIGN.CENTER

# Save presentation
prs.save(output_file)

print(f"\n✓ PowerPoint presentation saved to: {output_file}\n")
print("Layout:")
print("  A (top-left):     Before ComBat - Dataset")
print("  B (top-right):    After ComBat - Dataset")
print("  C (bottom-left):  Before ComBat - Trimester")
print("  D (bottom-right): After ComBat - Trimester\n")
print("All text boxes (title, labels, descriptions, caption) are editable in PowerPoint.\n")
print("==================================================================")
print("                    DONE                                          ")
print("==================================================================\n")
